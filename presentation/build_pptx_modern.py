#!/usr/bin/env python3
"""Clean modern presentation - Apple keynote inspired."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGURES = os.path.expanduser("~/Documents/research-workspace/dissertation/figures/advanced")
OUT = os.path.expanduser("~/Documents/research-workspace/dissertation/presentation/slides.pptx")

# Clean modern palette
BG_WHITE = RGBColor(255, 255, 255)
NEAR_BLACK = RGBColor(30, 30, 30)
DARK_TEXT = RGBColor(45, 45, 45)
SUBTITLE_GREY = RGBColor(120, 120, 120)
LIGHT_GREY_BG = RGBColor(245, 245, 247)
CARD_BORDER = RGBColor(230, 230, 232)
ACCENT_BLUE = RGBColor(0, 122, 255)  # Apple blue
ACCENT_GREEN = RGBColor(52, 199, 89)
ACCENT_ORANGE = RGBColor(255, 149, 0)
ACCENT_RED = RGBColor(255, 59, 48)
ACCENT_TEAL = RGBColor(90, 200, 250)
WHITE = RGBColor(255, 255, 255)
SOFT_BLUE_BG = RGBColor(240, 247, 255)
SOFT_GREEN_BG = RGBColor(240, 253, 244)
SOFT_ORANGE_BG = RGBColor(255, 248, 240)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL = 12

def add_slide(bg_color=BG_WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return s

def text(slide, left, top, width, height, txt="", size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Helvetica Neue"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def bullets(tf, items, size=16, color=DARK_TEXT, spacing=8, bold_indices=None):
    bold_indices = bold_indices or []
    for i, item in enumerate(items):
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Helvetica Neue"
        p.font.bold = i in bold_indices
        p.space_before = Pt(spacing)

def card(slide, left, top, width, height, fill=LIGHT_GREY_BG, border=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.05
    return shape

def divider(slide, left, top, width, color=CARD_BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def page_num(slide, num):
    text(slide, 12.6, 7.0, 0.6, 0.35, str(num), size=11, color=SUBTITLE_GREY, align=PP_ALIGN.RIGHT)

def accent_dot(slide, left, top, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(0.12), Inches(0.12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# =============================================
# SLIDE 1: Title
# =============================================
s = add_slide()
# Subtle accent line at top
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

text(s, 1.5, 2.0, 10.3, 1.5, "Identifying Depression\nThrough Speech",
     size=44, bold=True, color=NEAR_BLACK, align=PP_ALIGN.LEFT)
text(s, 1.5, 3.8, 10.3, 0.6, "A Comparative Analysis of Acoustic Features in Read and Spontaneous Speech",
     size=20, color=SUBTITLE_GREY, align=PP_ALIGN.LEFT)
divider(s, 1.5, 4.7, 2.0, ACCENT_BLUE)
text(s, 1.5, 5.0, 10.3, 0.4, "Nile", size=20, bold=True, color=DARK_TEXT)
text(s, 1.5, 5.5, 10.3, 0.8, "University of Glasgow  ·  School of Computing Science  ·  Level 4 Dissertation",
     size=14, color=SUBTITLE_GREY)
page_num(s, 1)

# =============================================
# SLIDE 2: Why This Matters
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 5, 0.6, "Why This Matters", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

# Problem card
card(s, 1.0, 1.5, 5.3, 4.8, LIGHT_GREY_BG)
text(s, 1.4, 1.7, 4.5, 0.5, "The Problem", size=20, bold=True, color=ACCENT_RED)

tf = text(s, 1.4, 2.3, 4.5, 3.8, "", size=16)
bullets(tf, [
    "280+ million people affected globally",
    "1 in 6 UK adults experience it weekly",
    "Diagnosis relies on subjective self-report",
    "Inter-rater reliability as low as 0.28 κ",
    "Significant access barriers to specialist care"
], size=16, spacing=10)

# Opportunity card
card(s, 7.0, 1.5, 5.3, 4.8, SOFT_BLUE_BG)
text(s, 7.4, 1.7, 4.5, 0.5, "The Opportunity", size=20, bold=True, color=ACCENT_BLUE)

card(s, 7.4, 2.5, 4.5, 3.2, WHITE, ACCENT_BLUE)
text(s, 7.6, 2.7, 4.1, 0.5, "🔊  Speech Analysis", size=20, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
divider(s, 7.8, 3.35, 3.7, CARD_BORDER)
tf = text(s, 7.6, 3.5, 4.1, 2.0, "", size=16)
bullets(tf, [
    "Non-invasive",
    "Can be collected remotely",
    "No specialist equipment needed",
    "Objective and repeatable"
], size=16, color=SUBTITLE_GREY, spacing=8)

page_num(s, 2)

# =============================================
# SLIDE 3: Research Gap
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "The Research Gap", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

# What exists
card(s, 1.0, 1.5, 5.3, 2.5, LIGHT_GREY_BG)
text(s, 1.4, 1.7, 4.5, 0.4, "What exists", size=18, bold=True, color=SUBTITLE_GREY)
text(s, 1.4, 2.2, 4.5, 0.8, "Deep learning models achieving 80%+ accuracy", size=17, color=DARK_TEXT)
card(s, 1.8, 3.0, 4.0, 0.7, WHITE, CARD_BORDER)
text(s, 2.0, 3.05, 3.6, 0.6, "\"Black Box\"  —  High accuracy, no explanation", size=14, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

# What's missing
card(s, 7.0, 1.5, 5.3, 2.5, SOFT_ORANGE_BG)
text(s, 7.4, 1.7, 4.5, 0.4, "What's missing", size=18, bold=True, color=ACCENT_ORANGE)
tf = text(s, 7.4, 2.2, 4.5, 1.5, "", size=17)
bullets(tf, [
    "Which features actually indicate depression?",
    "Does speech task type change what we find?"
], size=17, spacing=10)

# Research question
card(s, 1.0, 4.5, 11.3, 2.2, WHITE, ACCENT_BLUE)
accent_dot(s, 1.3, 4.75, ACCENT_BLUE)
text(s, 1.6, 4.65, 2.5, 0.4, "Research Question", size=14, bold=True, color=ACCENT_BLUE)
text(s, 1.4, 5.2, 10.5, 1.2,
     "Which acoustic features are most predictive of depression,\nand how do they differ between read and spontaneous speech?",
     size=22, color=DARK_TEXT, align=PP_ALIGN.CENTER)

page_num(s, 3)

# =============================================
# SLIDE 4: Approach Overview
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "Approach Overview", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

# Three pipeline cards
cards_data = [
    ("Data", ACCENT_BLUE, "ANDROIDS Corpus", ["118 Italian speakers", "Clinical psychiatric diagnoses", "Same participants, both tasks", "Reading + Interview"]),
    ("Features", ACCENT_GREEN, "eGeMAPS Feature Set", ["88 acoustic features", "Standardised & interpretable", "Prosody, spectral, timing", "Each feature has meaning"]),
    ("Models", ACCENT_ORANGE, "SVM + Random Forest", ["Interpretable by design", "Feature importance rankings", "Not deep learning", "Understanding over accuracy"]),
]

for i, (title, accent, subtitle, items) in enumerate(cards_data):
    x = 1.0 + i * 4.1
    card(s, x, 1.5, 3.6, 5.0, WHITE, CARD_BORDER)
    # Accent stripe at top of card
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.5), Inches(3.6), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    text(s, x + 0.3, 1.8, 3.0, 0.5, title, size=24, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    text(s, x + 0.3, 2.3, 3.0, 0.4, subtitle, size=14, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)
    divider(s, x + 0.4, 2.85, 2.8, CARD_BORDER)
    tf = text(s, x + 0.3, 3.0, 3.0, 3.0, "", size=15)
    bullets(tf, items, size=15, color=DARK_TEXT, spacing=8)

# Arrows
for i in range(2):
    x = 4.6 + i * 4.1
    text(s, x, 3.3, 0.5, 0.5, "→", size=28, bold=True, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

page_num(s, 4)

# =============================================
# SLIDE 5: Headline Result
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "The Headline Result", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

text(s, 1.0, 1.4, 12, 0.5, "Interview speech significantly outperforms reading",
     size=22, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Chart
fig_path = os.path.join(FIGURES, "accuracy_comparison.png")
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.8), Inches(2.1), Inches(6.0), Inches(4.5))

# Stats cards on right
card(s, 7.5, 2.1, 4.8, 1.4, LIGHT_GREY_BG)
text(s, 7.8, 2.2, 4.2, 0.4, "Reading", size=16, color=SUBTITLE_GREY)
text(s, 7.8, 2.55, 4.2, 0.6, "72.3%", size=36, bold=True, color=SUBTITLE_GREY)

card(s, 7.5, 3.7, 4.8, 1.4, SOFT_BLUE_BG, ACCENT_BLUE)
text(s, 7.8, 3.8, 4.2, 0.4, "Interview", size=16, color=ACCENT_BLUE)
text(s, 7.8, 4.15, 4.2, 0.6, "87.1%", size=36, bold=True, color=ACCENT_BLUE)

# Key metrics
card(s, 7.5, 5.4, 4.8, 1.5, WHITE, CARD_BORDER)
tf = text(s, 7.7, 5.5, 4.4, 1.3, "", size=15)
bullets(tf, [
    "↑  +15 percentage points",
    "📊  p = 0.0055 (statistically significant)",
    "🏆  Beats prior deep learning (81.6%)"
], size=15, spacing=6)

page_num(s, 5)

# =============================================
# SLIDE 6: Feature Divergence
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 11, 0.6, "Different Tasks, Different Features", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

text(s, 1.0, 1.4, 12, 0.4, "Of the top 10 features for each task, only 1 appears in both lists",
     size=19, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

# Reading card
card(s, 1.0, 2.2, 5.3, 2.8, LIGHT_GREY_BG)
accent_dot(s, 1.3, 2.45, SUBTITLE_GREY)
text(s, 1.6, 2.35, 4.5, 0.4, "Reading Task", size=20, bold=True, color=DARK_TEXT)
divider(s, 1.3, 2.9, 4.5, CARD_BORDER)
tf = text(s, 1.4, 3.0, 4.5, 1.8, "", size=16)
bullets(tf, [
    "Spectral slopes, MFCCs",
    "Static voice properties",
    "What the voice sounds like"
], size=16, spacing=8)

# Interview card
card(s, 7.0, 2.2, 5.3, 2.8, SOFT_BLUE_BG, ACCENT_BLUE)
accent_dot(s, 7.3, 2.45, ACCENT_BLUE)
text(s, 7.6, 2.35, 4.5, 0.4, "Interview Task", size=20, bold=True, color=ACCENT_BLUE)
divider(s, 7.3, 2.9, 4.5, ACCENT_BLUE)
tf = text(s, 7.4, 3.0, 4.5, 1.8, "", size=16)
bullets(tf, [
    "Variability measures, temporal dynamics",
    "Dynamic voice properties",
    "How the voice changes over time"
], size=16, spacing=8)

# Insight box
card(s, 1.0, 5.4, 11.3, 1.5, WHITE, ACCENT_ORANGE)
text(s, 1.4, 5.55, 10.5, 1.2,
     "These tasks don't give stronger or weaker versions of the same signal —\nthey provide access to qualitatively different information about depression.",
     size=17, color=DARK_TEXT, align=PP_ALIGN.CENTER)

page_num(s, 6)

# =============================================
# SLIDE 7: Variability Finding
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "The Variability Finding", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

# Chart
fig_path = os.path.join(FIGURES, "feature_importance_comparison.png")
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.2))

# Interview stats
card(s, 7.2, 1.5, 5.3, 2.5, SOFT_BLUE_BG, ACCENT_BLUE)
text(s, 7.5, 1.65, 4.7, 0.4, "Interview: 7/10 are variability", size=18, bold=True, color=ACCENT_BLUE)
divider(s, 7.5, 2.15, 4.5, ACCENT_BLUE)
tf = text(s, 7.5, 2.3, 4.7, 1.5, "", size=15)
bullets(tf, [
    "① Spectral flux variability",
    "② Hammarberg index variability",
    "③ Pause duration variability"
], size=15, spacing=6)

# Reading stats
card(s, 7.2, 4.3, 5.3, 0.8, LIGHT_GREY_BG)
text(s, 7.5, 4.4, 4.7, 0.5, "Reading: Only 2/10 are variability", size=16, bold=True, color=SUBTITLE_GREY)

# Key insight
card(s, 7.2, 5.4, 5.3, 1.5, SOFT_GREEN_BG, ACCENT_GREEN)
text(s, 7.5, 5.5, 4.7, 0.3, "Key Insight", size=13, bold=True, color=ACCENT_GREEN)
text(s, 7.5, 5.85, 4.7, 0.9,
     "Depression compresses dynamic range.\n\"Flat affect\" = reduced modulation\nacross multiple voice dimensions.",
     size=15, color=DARK_TEXT)

page_num(s, 7)

# =============================================
# SLIDE 8: Why Interview Works Better
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "Why Interview Works Better", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)
text(s, 1.0, 1.4, 12, 0.4, "Three mechanisms from clinical literature",
     size=18, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

mechanisms = [
    ("🧠", "Cognitive Load", ACCENT_BLUE, SOFT_BLUE_BG,
     "Spontaneous speech requires planning, retrieval, and monitoring",
     "Depression impairs executive function → pauses, hesitations"),
    ("❤️", "Emotional Engagement", ACCENT_RED, RGBColor(255, 245, 245),
     "Interview topics engage affect — mood, relationships, daily life",
     "Depression alters expression → only visible when emotions are relevant"),
    ("🎵", "Prosodic Freedom", ACCENT_GREEN, SOFT_GREEN_BG,
     "Conversation allows full control over intonation and rhythm",
     "Reading follows punctuation → suppresses natural modulation"),
]

for i, (icon, title, accent, bg, main, sub) in enumerate(mechanisms):
    x = 1.0 + i * 4.1
    card(s, x, 1.9, 3.6, 4.2, bg, accent)
    text(s, x, 2.1, 3.6, 0.6, icon, size=36, align=PP_ALIGN.CENTER)
    text(s, x + 0.3, 2.7, 3.0, 0.5, title, size=20, bold=True, color=accent, align=PP_ALIGN.CENTER)
    divider(s, x + 0.4, 3.3, 2.8, accent)
    text(s, x + 0.3, 3.5, 3.0, 1.2, main, size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    text(s, x + 0.3, 4.7, 3.0, 1.0, sub, size=12, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

# Bottom statement
card(s, 1.0, 6.4, 11.3, 0.7, LIGHT_GREY_BG)
text(s, 1.3, 6.45, 10.7, 0.6,
     "Interview captures all three channels. Reading suppresses them.",
     size=17, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

page_num(s, 8)

# =============================================
# SLIDE 9: Implications
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "Implications", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

implications = [
    ("Clinical Practice", ACCENT_BLUE, SOFT_BLUE_BG, [
        "Task selection > algorithm choice",
        "Use spontaneous speech for screening",
        "Measure variability, not just averages",
    ]),
    ("Research", ACCENT_GREEN, SOFT_GREEN_BG, [
        "Single-task studies miss markers",
        "Feature importance is task-contingent",
        "Don't generalise across tasks",
    ]),
    ("System Design", ACCENT_ORANGE, SOFT_ORANGE_BG, [
        "Elicit spontaneous speech",
        "Use open-ended questions",
        "Prioritise temporal dynamics",
    ]),
]

for i, (title, accent, bg, items) in enumerate(implications):
    x = 1.0 + i * 4.1
    card(s, x, 1.5, 3.6, 4.5, bg, accent)
    # Icon area
    icons = ["🩺", "🔬", "⚙️"]
    text(s, x, 1.7, 3.6, 0.5, icons[i], size=30, align=PP_ALIGN.CENTER)
    text(s, x + 0.3, 2.3, 3.0, 0.4, title, size=19, bold=True, color=accent, align=PP_ALIGN.CENTER)
    divider(s, x + 0.4, 2.8, 2.8, accent)
    tf = text(s, x + 0.3, 3.0, 3.0, 2.5, "", size=15)
    bullets(tf, items, size=15, spacing=10)

page_num(s, 9)

# =============================================
# SLIDE 10: Limitations
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "Limitations", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

# Most serious
card(s, 1.0, 1.5, 5.3, 3.5, RGBColor(255, 245, 245), ACCENT_RED)
accent_dot(s, 1.3, 1.75, ACCENT_RED)
text(s, 1.6, 1.65, 4.5, 0.4, "Most Serious", size=18, bold=True, color=ACCENT_RED)
divider(s, 1.3, 2.15, 4.5, RGBColor(255, 220, 220))

text(s, 1.4, 2.3, 4.5, 0.4, "Italian speech only", size=17, bold=True, color=DARK_TEXT)
text(s, 1.4, 2.7, 4.5, 0.5, "Features may not transfer to English or other languages — prosodic patterns can be language-dependent",
     size=14, color=SUBTITLE_GREY)

text(s, 1.4, 3.4, 4.5, 0.4, "Modest sample size", size=17, bold=True, color=DARK_TEXT)
text(s, 1.4, 3.8, 4.5, 0.5, "118 speakers — though results were stable across all 5 CV folds",
     size=14, color=SUBTITLE_GREY)

# Moderate
card(s, 7.0, 1.5, 5.3, 3.5, LIGHT_GREY_BG)
accent_dot(s, 7.3, 1.75, ACCENT_ORANGE)
text(s, 7.6, 1.65, 4.5, 0.4, "Moderate", size=18, bold=True, color=ACCENT_ORANGE)
divider(s, 7.3, 2.15, 4.5, CARD_BORDER)

text(s, 7.4, 2.3, 4.5, 0.4, "Binary labels only", size=17, bold=True, color=DARK_TEXT)
text(s, 7.4, 2.7, 4.5, 0.5, "No severity information — depressed or not, nothing in between",
     size=14, color=SUBTITLE_GREY)

text(s, 7.4, 3.4, 4.5, 0.4, "Gender imbalance (72% female)", size=17, bold=True, color=DARK_TEXT)
text(s, 7.4, 3.8, 4.5, 0.5, "Women misclassified slightly more often",
     size=14, color=SUBTITLE_GREY)

# Bottom takeaway
card(s, 1.0, 5.5, 11.3, 1.2, SOFT_BLUE_BG, ACCENT_BLUE)
text(s, 1.4, 5.65, 10.5, 0.9,
     "→  Cross-corpus validation across languages and populations is essential before clinical deployment",
     size=18, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

page_num(s, 10)

# =============================================
# SLIDE 11: Conclusion
# =============================================
s = add_slide()
text(s, 1.0, 0.5, 10, 0.6, "Conclusion", size=32, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 1.15, 1.5, ACCENT_BLUE)

# Quote card
card(s, 1.5, 1.5, 10.3, 1.6, SOFT_BLUE_BG, ACCENT_BLUE)
text(s, 1.9, 1.6, 9.5, 1.4,
     "\"Spontaneous speech captures depression markers that reading misses —\nspecifically, the reduced dynamic modulation that clinicians call 'flat affect'.\"",
     size=20, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

text(s, 1.0, 3.5, 10, 0.5, "Contributions", size=22, bold=True, color=NEAR_BLACK)
divider(s, 1.0, 4.05, 1.5, ACCENT_BLUE)

# 4 contribution cards
contribs = [
    ("01", "Task Effect", "+15 points, p = 0.0055"),
    ("02", "Task-Specific Features", "Only 1/10 overlap"),
    ("03", "Interpretable ≥ Deep Learning", "87.1% vs 81.6%"),
    ("04", "Variability is Key", "7/10 top interview features"),
]

for i, (num, title, detail) in enumerate(contribs):
    x = 1.0 + i * 3.05
    card(s, x, 4.3, 2.7, 2.5, WHITE, CARD_BORDER)
    text(s, x + 0.2, 4.4, 2.3, 0.5, num, size=28, bold=True, color=ACCENT_BLUE)
    text(s, x + 0.2, 5.0, 2.3, 0.5, title, size=15, bold=True, color=DARK_TEXT)
    text(s, x + 0.2, 5.5, 2.3, 0.6, detail, size=13, color=SUBTITLE_GREY)

page_num(s, 11)

# =============================================
# SLIDE 12: Questions
# =============================================
s = add_slide(LIGHT_GREY_BG)
# Accent line at top
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
stripe.fill.solid()
stripe.fill.fore_color.rgb = ACCENT_BLUE
stripe.line.fill.background()

text(s, 1.0, 2.2, 11.3, 0.8, "Thank You", size=48, bold=True, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
text(s, 1.0, 3.2, 11.3, 0.5, "Questions?", size=28, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

divider(s, 5.5, 4.0, 2.3, ACCENT_BLUE)

text(s, 1.0, 4.3, 11.3, 0.5, "Nile", size=22, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
text(s, 1.0, 4.8, 11.3, 0.4, "University of Glasgow  ·  School of Computing Science",
     size=14, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

# Key numbers card
card(s, 3.0, 5.5, 7.3, 1.3, WHITE, CARD_BORDER)
text(s, 3.3, 5.6, 6.7, 1.1,
     "87.1% interview  vs  72.3% reading   ·   p = 0.0055\n1/10 feature overlap   ·   7/10 interview features are variability",
     size=14, color=SUBTITLE_GREY, align=PP_ALIGN.CENTER)

page_num(s, 12)

prs.save(OUT)
print(f"✅ Saved to {OUT}")
