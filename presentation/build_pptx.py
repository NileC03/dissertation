#!/usr/bin/env python3
"""Convert dissertation LaTeX Beamer slides to PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGURES = os.path.expanduser("~/Documents/research-workspace/dissertation/figures/advanced")
OUT = os.path.expanduser("~/Documents/research-workspace/dissertation/presentation/slides.pptx")

# Colors
UOFG_BLUE = RGBColor(0, 78, 124)
ACCENT_BLUE = RGBColor(0, 119, 182)
DARK_GREY = RGBColor(73, 80, 87)
LIGHT_GREY = RGBColor(248, 249, 250)
SUCCESS_GREEN = RGBColor(25, 135, 84)
WARNING_ORANGE = RGBColor(255, 193, 7)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text="", font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullet_list(tf, items, font_size=16, color=BLACK, bold_items=None):
    bold_items = bold_items or []
    for i, item in enumerate(items):
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = i in bold_items
        p.level = 0
        p.space_before = Pt(4)

def add_rounded_box(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def slide_number(slide, num, total):
    add_textbox(slide, 12.5, 7.0, 0.7, 0.4, f"{num}/{total}", font_size=10, color=DARK_GREY, alignment=PP_ALIGN.RIGHT)

TOTAL = 12

# ---- SLIDE 1: Title ----
s = add_slide()
# Background
add_textbox(s, 1.5, 1.2, 10.3, 1.2, "Identifying Depression Through Speech",
            font_size=36, bold=True, color=UOFG_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 2.5, 10.3, 0.8,
            "A Comparative Analysis of Acoustic Features\nin Read and Spontaneous Speech",
            font_size=20, color=DARK_GREY, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 3.8, 10.3, 0.5, "Nile",
            font_size=22, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 4.4, 10.3, 0.8,
            "University of Glasgow\nSchool of Computing Science\nLevel 4 Dissertation",
            font_size=14, color=DARK_GREY, alignment=PP_ALIGN.CENTER)
slide_number(s, 1, TOTAL)

# ---- SLIDE 2: Why This Matters ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "Why This Matters", font_size=30, bold=True, color=UOFG_BLUE)

add_textbox(s, 0.8, 1.2, 5.5, 0.5, "⚠ The Problem", font_size=22, bold=True, color=UOFG_BLUE)
tf = add_textbox(s, 0.8, 1.8, 5.5, 4.0, "", font_size=18)
add_bullet_list(tf, [
    "280+ million affected globally",
    "1 in 6 UK adults weekly",
    "Diagnosis relies on self-report",
    "Inter-rater reliability as low as 0.28 kappa",
    "Access barriers to specialist care"
], font_size=18, bold_items=[0, 1])

add_textbox(s, 7.0, 1.2, 5.5, 0.5, "💡 The Opportunity", font_size=22, bold=True, color=SUCCESS_GREEN)
box = add_rounded_box(s, 7.3, 2.0, 4.8, 2.5, LIGHT_GREY, ACCENT_BLUE)
add_textbox(s, 7.5, 2.2, 4.4, 2.0,
            "🔊 Speech\n\nNon-invasive\nRemote collection\nNo specialist equipment",
            font_size=18, color=BLACK, alignment=PP_ALIGN.CENTER)
slide_number(s, 2, TOTAL)

# ---- SLIDE 3: Research Gap ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "The Research Gap", font_size=30, bold=True, color=UOFG_BLUE)

add_textbox(s, 0.8, 1.2, 5.5, 0.5, "✓ What exists", font_size=22, bold=True, color=DARK_GREY)
add_textbox(s, 0.8, 1.8, 5.5, 1.0, "Deep learning models achieving 80%+ accuracy", font_size=18, color=BLACK)
box = add_rounded_box(s, 1.2, 3.0, 4.5, 1.5, RGBColor(230, 230, 230), DARK_GREY)
add_textbox(s, 1.4, 3.1, 4.1, 1.3, '"Black Box"\n\nHigh accuracy\nNo explanation',
            font_size=16, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 7.0, 1.2, 5.5, 0.5, "❓ What's missing", font_size=22, bold=True, color=ACCENT_BLUE)
tf = add_textbox(s, 7.0, 1.8, 5.5, 1.5, "", font_size=18)
add_bullet_list(tf, [
    "Which features indicate depression?",
    "Does speech task type matter?"
], font_size=18)

box = add_rounded_box(s, 7.0, 3.5, 5.5, 2.0, LIGHT_GREY, UOFG_BLUE)
add_textbox(s, 7.2, 3.6, 5.1, 0.4, "Research Question", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
# Hacky: put title bg
title_bg = add_rounded_box(s, 7.0, 3.5, 5.5, 0.5, UOFG_BLUE)
add_textbox(s, 7.2, 3.55, 5.1, 0.4, "Research Question", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(s, 7.2, 4.1, 5.1, 1.2,
            "Which acoustic features are most predictive of depression, and how do they differ between read and spontaneous speech?",
            font_size=16, color=DARK_GREY, alignment=PP_ALIGN.LEFT)
slide_number(s, 3, TOTAL)

# ---- SLIDE 4: Approach Overview ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "Approach Overview", font_size=30, bold=True, color=UOFG_BLUE)

# Pipeline boxes
for i, (title, details, desc) in enumerate([
    ("Data", "ANDROIDS\n118 speakers\nBoth tasks", "Clinical diagnoses\nSame participants do both tasks"),
    ("Features", "eGeMAPS\n88 interpretable\nacoustic features", "Prosody, spectral, voice\nquality, timing\nEach feature = meaning"),
    ("Models", "SVM + RF\nInterpretable\nNot \"black box\"", "Deliberately interpretable\nNot deep learning"),
]):
    x = 1.5 + i * 4.0
    box = add_rounded_box(s, x, 1.3, 3.2, 2.2, LIGHT_GREY, UOFG_BLUE)
    add_textbox(s, x + 0.1, 1.4, 3.0, 0.4, title, font_size=20, bold=True, color=UOFG_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, 1.9, 3.0, 1.5, details, font_size=15, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + 0.1, 4.0, 3.0, 1.5, desc, font_size=13, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

# Arrows between boxes
for i in range(2):
    x = 4.7 + i * 4.0
    add_textbox(s, x, 2.0, 0.8, 0.6, "→", font_size=32, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

slide_number(s, 4, TOTAL)

# ---- SLIDE 5: Headline Result ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "The Headline Result", font_size=30, bold=True, color=UOFG_BLUE)
add_textbox(s, 0.5, 1.0, 12, 0.5, "Interview speech significantly outperforms reading",
            font_size=24, bold=True, color=UOFG_BLUE, alignment=PP_ALIGN.CENTER)

fig_path = os.path.join(FIGURES, "accuracy_comparison.png")
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0))

# Results on right
add_textbox(s, 7.8, 2.0, 4.5, 0.4, "Task                    Accuracy", font_size=16, bold=True, color=BLACK)
add_textbox(s, 7.8, 2.5, 4.5, 0.4, "Reading                72.3%", font_size=16, color=BLACK)
add_textbox(s, 7.8, 2.9, 4.5, 0.4, "Interview              87.1%", font_size=18, bold=True, color=UOFG_BLUE)

tf = add_textbox(s, 7.8, 3.8, 4.5, 3.0, "", font_size=18)
add_bullet_list(tf, [
    "↑ +15 points difference",
    "📈 p = 0.0055 significant",
    "🏆 Beats prior DL (81.6%)"
], font_size=17, bold_items=[0])
slide_number(s, 5, TOTAL)

# ---- SLIDE 6: Feature Divergence ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "But Here's What's Interesting...", font_size=30, bold=True, color=UOFG_BLUE)
add_textbox(s, 0.5, 1.0, 12, 0.5, "Different features predict depression in each task",
            font_size=22, bold=True, color=UOFG_BLUE, alignment=PP_ALIGN.CENTER)

box = add_rounded_box(s, 1.0, 1.8, 5.0, 2.0, LIGHT_GREY, UOFG_BLUE)
title_bg = add_rounded_box(s, 1.0, 1.8, 5.0, 0.5, UOFG_BLUE)
add_textbox(s, 1.2, 1.85, 4.6, 0.4, "Reading Task", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(s, 1.2, 2.4, 4.6, 1.2, "Spectral slopes\nMFCCs\n\nStatic properties",
            font_size=16, color=BLACK, alignment=PP_ALIGN.CENTER)

box = add_rounded_box(s, 7.3, 1.8, 5.0, 2.0, LIGHT_GREY, UOFG_BLUE)
title_bg = add_rounded_box(s, 7.3, 1.8, 5.0, 0.5, UOFG_BLUE)
add_textbox(s, 7.5, 1.85, 4.6, 0.4, "Interview Task", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(s, 7.5, 2.4, 4.6, 1.2, "Variability measures\nTemporal dynamics\n\nDynamic properties",
            font_size=16, color=BLACK, alignment=PP_ALIGN.CENTER)

box = add_rounded_box(s, 1.5, 4.5, 10.3, 1.5, RGBColor(255, 248, 230), WARNING_ORANGE)
add_textbox(s, 1.7, 4.6, 9.9, 1.3,
            "⚠ Of top 10 features: Only 1 overlaps\n\nThis isn't just \"interview is better\" — they capture qualitatively different information",
            font_size=17, color=BLACK, alignment=PP_ALIGN.CENTER)
slide_number(s, 6, TOTAL)

# ---- SLIDE 7: Variability Finding ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "The Variability Finding", font_size=30, bold=True, color=UOFG_BLUE)

fig_path = os.path.join(FIGURES, "feature_importance_comparison.png")
if os.path.exists(fig_path):
    s.shapes.add_picture(fig_path, Inches(0.5), Inches(1.2), Inches(6.5), Inches(5.5))

add_textbox(s, 7.5, 1.2, 5.0, 0.5, "Interview: 7/10 are variability", font_size=20, bold=True, color=UOFG_BLUE)
tf = add_textbox(s, 7.5, 1.8, 5.0, 2.0, "", font_size=17)
add_bullet_list(tf, [
    "1. Spectral flux variability",
    "2. Hammarberg index variability",
    "3. Pause duration variability",
], font_size=17)

add_textbox(s, 7.5, 3.5, 5.0, 0.5, "Reading: Only 2/10", font_size=18, bold=True, color=DARK_GREY)

box = add_rounded_box(s, 7.5, 4.3, 5.0, 2.0, RGBColor(230, 250, 240), SUCCESS_GREEN)
add_textbox(s, 7.7, 4.4, 4.6, 1.8,
            "Key insight:\n\nDepression compresses dynamic range\n\n\"Flat affect\" = reduced modulation",
            font_size=16, bold=False, color=BLACK, alignment=PP_ALIGN.CENTER)
slide_number(s, 7, TOTAL)

# ---- SLIDE 8: Why Interview Works Better ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "Why Interview Works Better", font_size=30, bold=True, color=UOFG_BLUE)
add_textbox(s, 0.5, 1.0, 12, 0.5, "Three mechanisms from clinical literature",
            font_size=22, bold=True, color=UOFG_BLUE, alignment=PP_ALIGN.CENTER)

for i, (icon, title, main, sub) in enumerate([
    ("🧠", "Cognitive Load",
     "Spontaneous speech requires planning, retrieval, monitoring",
     "Depression impairs executive function → pauses, hesitations"),
    ("❤️", "Emotional Engagement",
     "Interview topics engage affect (mood, relationships)",
     "Depression alters expression → only visible when emotions relevant"),
    ("〰️", "Prosodic Freedom",
     "Reading follows punctuation; conversation allows natural modulation",
     "Depression compresses this natural range"),
]):
    x = 0.8 + i * 4.2
    add_textbox(s, x, 1.7, 3.5, 0.5, icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, 2.3, 3.5, 0.5, title, font_size=20, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, 2.9, 3.5, 1.5, main, font_size=15, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, 4.3, 3.5, 1.5, sub, font_size=13, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 0.5, 6.0, 12, 0.5, "Interview captures all three channels. Reading suppresses them.",
            font_size=16, color=DARK_GREY, alignment=PP_ALIGN.CENTER)
slide_number(s, 8, TOTAL)

# ---- SLIDE 9: Implications ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "Implications", font_size=30, bold=True, color=UOFG_BLUE)

for i, (title, items) in enumerate([
    ("🩺 Clinical Practice", [
        "Task selection > algorithm choice",
        "Spontaneous > reading for screening",
        "Measure variability, not just means",
    ]),
    ("🔬 Research", [
        "Single-task studies miss markers",
        "Feature importance is task-contingent",
        "\"MFCCs are best\" may only hold for reading",
    ]),
    ("⚙️ System Design", [
        "Elicit spontaneous speech",
        "Use open-ended questions",
        "Prioritise temporal dynamics",
    ]),
]):
    x = 0.8 + i * 4.2
    add_textbox(s, x, 1.2, 3.5, 0.5, title, font_size=20, bold=True, color=UOFG_BLUE)
    tf = add_textbox(s, x, 1.8, 3.5, 4.0, "", font_size=16)
    add_bullet_list(tf, items, font_size=16)

slide_number(s, 9, TOTAL)

# ---- SLIDE 10: Limitations ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "Limitations", font_size=30, bold=True, color=UOFG_BLUE)

add_textbox(s, 0.8, 1.2, 5.5, 0.5, "⚠ Most Serious", font_size=22, bold=True, color=WARNING_ORANGE)
tf = add_textbox(s, 0.8, 1.8, 5.5, 3.0, "", font_size=17)
add_bullet_list(tf, [
    "Italian speech only",
    "  → Features may not transfer to English",
    "",
    "Modest sample size (118 speakers)",
], font_size=17, bold_items=[0, 3])

add_textbox(s, 7.0, 1.2, 5.5, 0.5, "ℹ Moderate", font_size=22, bold=True, color=DARK_GREY)
tf = add_textbox(s, 7.0, 1.8, 5.5, 3.0, "", font_size=17)
add_bullet_list(tf, [
    "Binary labels only",
    "  → No severity information",
    "",
    "Gender imbalance (72% female)",
    "  → Women misclassified slightly more",
], font_size=17, bold_items=[0, 3])

box = add_rounded_box(s, 1.5, 5.2, 10.3, 1.2, RGBColor(230, 240, 255), ACCENT_BLUE)
add_textbox(s, 1.7, 5.3, 9.9, 1.0,
            "→ Cross-corpus validation essential before clinical deployment",
            font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
slide_number(s, 10, TOTAL)

# ---- SLIDE 11: Conclusion ----
s = add_slide()
add_textbox(s, 0.5, 0.3, 12, 0.7, "Conclusion", font_size=30, bold=True, color=UOFG_BLUE)

box = add_rounded_box(s, 1.0, 1.2, 11.3, 1.8, RGBColor(235, 242, 248), UOFG_BLUE)
add_textbox(s, 1.3, 1.3, 10.7, 1.6,
            "\"Spontaneous speech captures depression markers that reading misses —\nspecifically, the reduced dynamic modulation that clinicians call 'flat affect'.\"",
            font_size=20, color=UOFG_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(s, 0.8, 3.3, 12, 0.5, "Contributions:", font_size=22, bold=True, color=UOFG_BLUE)

for i, (contrib, detail) in enumerate([
    ("1. Quantified the task effect", "+15 points, p = 0.0055"),
    ("2. Showed features are task-specific", "Only 1/10 overlap"),
    ("3. Interpretable methods match DL", "87.1% vs 81.6% prior work"),
    ("4. Identified key signal: variability", "7/10 top interview features"),
]):
    col = 0 if i < 2 else 1
    row = i % 2
    x = 0.8 + col * 6.2
    y = 4.0 + row * 1.2
    add_textbox(s, x, y, 5.5, 0.5, contrib, font_size=18, bold=True, color=BLACK)
    add_textbox(s, x, y + 0.4, 5.5, 0.4, detail, font_size=14, color=DARK_GREY)

slide_number(s, 11, TOTAL)

# ---- SLIDE 12: Questions ----
s = add_slide()
add_textbox(s, 0.5, 1.5, 12, 1.0, "💬", font_size=60, alignment=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 2.8, 12, 0.8, "Thank you", font_size=36, bold=True, color=UOFG_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 3.8, 12, 0.5, "Nile", font_size=22, color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 4.3, 12, 0.4, "University of Glasgow", font_size=16, color=DARK_GREY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 5.5, 10.3, 1.0,
            "Key numbers: 87.1% interview vs 72.3% reading (p = 0.0055)\n1/10 feature overlap  •  7/10 interview features are variability",
            font_size=14, color=DARK_GREY, alignment=PP_ALIGN.CENTER)
slide_number(s, 12, TOTAL)

prs.save(OUT)
print(f"Saved to {OUT}")
